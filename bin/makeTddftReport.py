import numpy as np
import matplotlib.pyplot as plt
import pyscf
from pyscf.geomopt.geometric_solver import optimize
from pyscf import gto, scf, dft, tddft
from rdkit import Chem
from rdkit.Chem import rdDepictor
from rdkit.Chem.Draw import rdMolDraw2D
from rdkit.Chem.Draw.MolDrawing import DrawingOptions

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

import subprocess
import s00Utils
import math,os,sys

def showMapPeak(xx,yy):
    m=0.0
    for x,y in zip(xx,yy):
        if y > m:
            m=y
            b=x

    return '{:.1f}'.format(b)+":"+'{:.2f}'.format(m)

def showPeak(band,strength,n):
    m=0.0
    for i,s in enumerate(strength):
        if s > m:
            m=s
            b=band[i]

    return '{:.1f}'.format(b)+":"+'{:.2e}'.format(m)

def gaussBand(x, band, strength, stdev):
  constant=0.0003
  bandshape = constant * (strength/(1.0/stdev))*np.exp(-(((1.0/x)-(1.0/band))**2/(1.0/stdev)**2))
  return bandshape

def peak(xs):
    m=0.0
    for x in xs:
        if m<x:
            m=x
    return f':{m:.2f}'

def norm(ys,tp,tpp):
    ret=[]
    m=0
    if type(tp)==type(False):
        for y in ys:
            if m<y:
                m=y
    else:
        rx=[300,400]
        if tpp=='type4':
            rx=[220,300]
        print("------",rx)
        for x,y in zip(tp,ys):
            if x<rx[0] or x>rx[1]:
                continue
            if m<y:
                m=y

    for y in ys:
        ret.append(y/m)
    return ret

def putCell(tbl,i,j,txt,size=12):
    tbl.cell(i,j).text=txt
    tbl.cell(i,j).text_frame.paragraphs[0].font.size=Pt(size)


def highPeakExist(w,f,thr,rl):

    if rl is not False:
        return agent(w,f,thr,rl)

    for i in w:
        if i > float(thr):
            return True
    return False

def agent(w,f,thr,rl):
    ff=[];fm=0.0
    for i in f:
        if fm < i:
            fm=i
    for i in f:
        ff.append(i/fm)

    for i,j in zip(w,ff):
        if i > float(thr) and j > rl:
            return True
    return False

def strokePage(l):
    print("locals",l)
    ax=l['ax'];ym=l['ym'];xs=l['xs']
    plt=l['plt'];fno=l['fno']
    tables=l['tables'];all_tables=l['all_tables']
    labels=l['labels'];all_labels=l['all_labels']
    peaks=l['peaks'];all_peaks=l['all_peaks']
    smiles=l['smiles']

    ax.set_ylim([0,ym])
    ax.set_xlim([int(xs[0]),int(xs[1])])
    ax.legend()
    plt.xlabel("wavelength(nm)")
    plt.ylabel(ylabel)
    s00Utils.saveAs(plt,fn=fno+".png")

    drawer = rdMolDraw2D.MolDraw2DSVG(600,750,300,150)
    drawer.SetLineWidth(2)
    drawer.SetFontSize(1.0)
    all_tables.append(tables)
    all_labels.append(labels)
    all_peaks.append(peaks)
#
    mols=[];ls=[];ss=[]
    for s,l in zip(smiles,labels):
        if s not in ss:
            mol=Chem.MolFromSmiles(s)
            mols.append(mol)
            ss.append(s)
            ls.append(l)

    drawer.DrawMolecules(mols,legends=ls)
    drawer.FinishDrawing()
    svg = drawer.GetDrawingText().replace('svg:', '')

    print("fno:"+fno)
    with open(fno+"s.svg",'w') as f:
        f.write(svg)

###############################
###### main
###############################
if len(sys.argv)<2:
    print("version 2.1")
    print("please specify parameter file")
    print(sys.argv[0],"options")
    print("-n|-N : normalized by each peak value, -N : describe peak values in labels")
    print("-s : normalized by each peak value existing between 300-400nm")
    print("-up threshold : plot tdds where the peak higher than threshold exists")
    print("-r 0.1 : the peak height is higher than 10%.")
    print("-num : specify the number of lines per graph");
    print("-no_table : no table")

    exit()

normal=False;table='homo';N_PER=10
threshold=False
relative=False
#sep=False
#if len(sys.argv)>2:
#    sep=True

if len(sys.argv)>2:
    for i in range(2,len(sys.argv)):
        if sys.argv[i]=='-n':
            normal='type1'
        elif sys.argv[i]=='-N':
            normal='type2'
        elif sys.argv[i]=='-S':
            normal='type4'
        elif sys.argv[i]=='-h':
            table='homo'
        elif sys.argv[i]=='-up':
            threshold=int(sys.argv[i+1])
        elif sys.argv[i]=='-r':
            relative=float(sys.argv[i+1])
        elif sys.argv[i]=='-num':
            N_PER=int(sys.argv[i+1])
        elif sys.argv[i]=='-no_table':
            table=False

with open(sys.argv[1],"r") as fp:
    fns=fp.readline().split(":")[1].split("\n")[0].split(',')
    if fns[-1]=='':
        fns=fns[:-1]
    fni=fp.readline().split(':')[1].split("\n")[0]
    ops=fp.readline().split(':')[1].split("\n")[0]
    xs=ops.split("]")[0].split("[")[1].split(',')
    ym=float(ops.split("]")[1].split(',')[1])
    stdev=float(ops.split("]")[1].split(',')[2])

    x=np.linspace(int(xs[0]),int(xs[1]),int(xs[2]))

    ylabel="Absorption (a.u.)"

    if normal!=False:
        ylabel="Normarized absorption"
        ym=1.1

loop=len(fns)
print("<--------",fns)

n=1
c=0
new=True

wk=os.getcwd().split('wsl')[1].split('/')

fni='uv'+"_".join(wk)+"_"+fni

fpo=open(fni+".tsv","w")

all_files=[];all_tables=[];all_labels=[];all_peaks=[]
peaks=[]
for fn in fns:
    if fn=='':
        continue
    w,f,s,h=s00Utils.getNP(fn)
    if threshold is not False:
        if not highPeakExist(w,f,threshold,relative):
            continue
    c+=1

    if len(w)<1:
        continue
    if new:
        fno=fni+"_"+str(n)
        all_files.append(fno)
        smiles=[];labels=[];tables=[];ids=[];peaks=[]
        Figure, ax=plt.subplots()
        new=False

    tables.append(h)

    print("The first peak ",showPeak(w,f,1))
#    if sep:
#        fn=fn.split('_')[1]

    composite=s00Utils.getCP(x,w,f,stdev)
    print("The first peak ",showMapPeak(x,composite))
    fnx=fn
    if normal=='type2':
        fnx=fn+peak(composite)

    smiles.append(s)
    labels.append(fn)

    if normal!=False:
        if normal=='type1' or normal=="type2":
            composite=norm(composite,False,False)
        else:
            composite=norm(composite,x,normal)

    xx = ', '.join([f'{num:.1f}' for num in x])
    yy = ', '.join([f'{num:.3f}' for num in composite])
    fpo.write(fn.split(".")[0]+'\t'+s+'\t'+xx+'\t'+yy+"\n")

    ax.plot(x,composite,label=fnx)
    peaks.append(s00Utils.getPeaks(x,composite))
    print("peaks",peaks)

    if c%N_PER==0:
        print(fno)
        new=True
        strokePage(locals())
        n=n+1

if not new:
    strokePage(locals())

page=Presentation()
sx=Inches(8.27);sy=Inches(11.69)
page.slide_height=sx;page.slide_width =sy

blank_slide_layout = page.slide_layouts[6]
slide = page.slides.add_slide(blank_slide_layout)

print(all_files)

n=0
for fno in all_files:
    if n!=0:
        slide = page.slides.add_slide(blank_slide_layout)
    n+=1
#    if not os.path.exists(fno+"s.png"):
    subprocess.run(['convert',fno+"s.svg",fno+"s.png"])
    slide.shapes.add_picture(fno+".png",left=sx*0.05,top=sy*0.25,width=sx*0.7)
    slide.shapes.add_picture(fno+"s.png",left=sx*0.75,top=sy*0.10,width=sx*0.6)

    if table:
        if table=='homo':
            yn=len(all_tables[n-1])
#            yt=Cm(15)-Cm(0.6)*(yn+1)
            yt=Cm(17)-Cm(0.6)*(yn+1)
            tbl=slide.shapes.add_table(yn+1,5,Cm(16),yt,Cm(12),Cm(0.6)*(yn+1)).table
            putCell(tbl,0,1,'HOMO (eV)');putCell(tbl,0,2,'LUMO (eV)');putCell(tbl,0,3,'gap (eV)')
            putCell(tbl,0,4,'λmax(nm)')

            for i,l in enumerate(all_tables[n-1]):
                putCell(tbl,i+1,0,all_labels[n-1][i])
                putCell(tbl,i+1,1,"{:.4f}".format(l[0]))
                putCell(tbl,i+1,2,"{:.4f}".format(l[1]))
                putCell(tbl,i+1,3,"{:.4f}".format(l[1]-l[0]))
#                putCell(tbl,i+1,4,all_peaks[n-1][i])
                putCell(tbl,i+1,4,all_peaks[n-1][i])

fno=fni+".pptx"
page.save(fno)
fpo.close()
